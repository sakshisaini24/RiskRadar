"""Generate a single 16:9 business-impact slide (import into Google Slides via .pptx)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "RiskRadar_Business_Impact.pptx"

# Matches frontend ROI defaults (page.tsx)
MONTHLY_VOLUME = 500
POSITIVE_RATE = 0.35
CONSERVATIVE_RECALL = 0.77
AVG_CLAIM = 289_000
ESCALATION_MULT = 3.5
MINUTES_SAVED = 37
ADJUSTER_HOURLY = 55

BG = RGBColor(6, 78, 59)         # emerald-900
PANEL = RGBColor(6, 95, 70)      # emerald-800
ACCENT = RGBColor(251, 191, 36)  # amber-300
MUTED = RGBColor(167, 243, 208)  # emerald-200
WHITE = RGBColor(255, 255, 255)


def _roi():
    esc = MONTHLY_VOLUME * POSITIVE_RATE
    early = esc * (CONSERVATIVE_RECALL - POSITIVE_RATE)
    cost_miss = AVG_CLAIM * (ESCALATION_MULT - 1)
    monthly_esc = early * cost_miss
    hours = esc * MINUTES_SAVED / 60
    labor_mo = hours * ADJUSTER_HOURLY
    return {
        "early_per_mo": round(early, 1),
        "hours_per_mo": round(hours),
        "labor_annual_k": round(labor_mo * 12 / 1000),
        "annual_esc_m": monthly_esc * 12 / 1_000_000,
        "cost_per_miss_k": round(cost_miss / 1000),
        "lift_pct": round((CONSERVATIVE_RECALL / POSITIVE_RATE - 1) * 100),
    }


def _box(slide, left, top, width, height, fill=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def _text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    return tb


def build():
    r = _roi()
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _box(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, BG)

    _text(
        slide, Inches(0.55), Inches(0.3), Inches(12), Inches(0.75),
        "Quantifiable business impact",
        size=34, bold=True,
    )
    _text(
        slide, Inches(0.55), Inches(1.0), Inches(11),
        Inches(0.35),
        f"Portfolio scenario: {MONTHLY_VOLUME:,} claims/mo · ${AVG_CLAIM // 1000}K avg · "
        f"{int(CONSERVATIVE_RECALL * 100)}% conservative recall (holdout-validated)",
        size=12, color=MUTED,
    )

    cards = [
        (f"{r['hours_per_mo']}", "Adjuster hours returned / month", f"37 min saved × deep reviews (45→8 min workflow)"),
        (f"{r['early_per_mo']}", "More escalations caught early / month", f"vs random triage at {int(POSITIVE_RATE * 100)}% baseline recall"),
        (f"{r['lift_pct']}%", "Higher catch rate", f"{int(CONSERVATIVE_RECALL * 100)}% model recall vs {int(POSITIVE_RATE * 100)}% random queue"),
        (f"${r['labor_annual_k']}K", "Annual labor savings", f"At ${ADJUSTER_HOURLY}/hr · capacity you can redeploy"),
    ]
    x = 0.55
    for big, title, sub in cards:
        _box(slide, Inches(x), Inches(1.55), Inches(3.05), Inches(2.35), PANEL)
        _text(slide, Inches(x + 0.2), Inches(1.75), Inches(2.7), Inches(0.85), big, size=40, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        _text(slide, Inches(x + 0.15), Inches(2.55), Inches(2.75), Inches(0.55), title, size=12, bold=True, align=PP_ALIGN.CENTER)
        _text(slide, Inches(x + 0.15), Inches(3.15), Inches(2.75), Inches(0.7), sub, size=9, color=MUTED, align=PP_ALIGN.CENTER)
        x += 3.2

    _box(slide, Inches(0.55), Inches(4.15), Inches(12.2), Inches(1.55), PANEL)
    _text(
        slide, Inches(0.75), Inches(4.3), Inches(5.5), Inches(0.35),
        "ESCALATION LOSS AVOIDANCE (MODELED)",
        size=10, bold=True, color=ACCENT,
    )
    _text(
        slide, Inches(0.75), Inches(4.65), Inches(4.5), Inches(0.9),
        f"${r['annual_esc_m']:.1f}M",
        size=44, bold=True,
    )
    _text(
        slide, Inches(0.75), Inches(5.35), Inches(5.2), Inches(0.35),
        "annual scenario at 3.5× escalated cost multiplier",
        size=11, color=MUTED,
    )
    _text(
        slide, Inches(5.8), Inches(4.35), Inches(6.5), Inches(1.2),
        f"Each early catch avoids ~${r['cost_per_miss_k']:,}K incremental legal + settlement drag "
        f"(avg claim × {ESCALATION_MULT - 1:.1f}×). Live ROI sliders in the RiskRadar dashboard — "
        f"judges can stress-test assumptions in the demo.",
        size=11, color=MUTED,
    )

    _text(
        slide, Inches(0.55), Inches(6.0), Inches(12.2), Inches(0.55),
        "Bottom line: faster triage · more escalations surfaced before legal fees compound · "
        "ROI tied to real holdout metrics, not gut feel.",
        size=13, bold=True, align=PP_ALIGN.CENTER,
    )
    _text(
        slide, Inches(0.55), Inches(6.65), Inches(12.2), Inches(0.35),
        "Assumptions match in-app ROI calculator · Planning mode · conservative 77% recall",
        size=9, color=MUTED, align=PP_ALIGN.CENTER,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")
    print(f"ROI: {r}")


if __name__ == "__main__":
    build()
