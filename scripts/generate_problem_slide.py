"""Generate a single 16:9 problem-statement slide (import into Google Slides via .pptx)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "RiskRadar_Problem_Statement.pptx"

# Palette
BG = RGBColor(15, 23, 42)       # slate-900
ACCENT = RGBColor(251, 146, 60)  # orange-400
MUTED = RGBColor(148, 163, 184)  # slate-400
WHITE = RGBColor(248, 250, 252)
BLUE = RGBColor(56, 189, 248)    # sky-400
RED = RGBColor(248, 113, 113)     # red-400


def _box(slide, left, top, width, height, fill=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)  # rectangle
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def _text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    return tb


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

    # Background
    _box(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, BG)

    # Accent bar top
    _box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT)

    # Title block
    _text(
        slide,
        Inches(0.6),
        Inches(0.35),
        Inches(12),
        Inches(0.9),
        "When escalation is obvious, it's already too late",
        size=36,
        bold=True,
        color=WHITE,
    )
    _text(
        slide,
        Inches(0.6),
        Inches(1.15),
        Inches(10),
        Inches(0.4),
        "Insurance claims  ·  legal risk  ·  early intervention",
        size=14,
        color=MUTED,
    )

    # Left column — pain points
    _text(slide, Inches(0.6), Inches(1.75), Inches(5.8), Inches(0.35), "THE PAIN", size=11, bold=True, color=ACCENT)
    bullets = [
        ("Specialists guess", "Gut feel + manual review of voluminous notes and structured data"),
        ("Reactive, not predictive", "Red flags appear after the low-cost resolution window closes"),
        ("Two costs", "Legal fees can exceed claim value · ~45 min/case on manual research"),
    ]
    y = 2.15
    for i, (head, body) in enumerate(bullets, 1):
        _box(slide, Inches(0.6), Inches(y), Inches(0.35), Inches(0.35), ACCENT if i == 3 else BLUE)
        num = slide.shapes[-1]
        tf = num.text_frame
        tf.paragraphs[0].text = str(i)
        tf.paragraphs[0].font.size = Pt(14)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = BG
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE

        _text(slide, Inches(1.05), Inches(y - 0.02), Inches(5.3), Inches(0.35), head, size=16, bold=True)
        _text(slide, Inches(1.05), Inches(y + 0.32), Inches(5.3), Inches(0.55), body, size=13, color=MUTED)
        y += 1.35

    _text(
        slide,
        Inches(0.6),
        Inches(5.35),
        Inches(5.8),
        Inches(0.9),
        "We need intelligence that reads between the lines—before cases hit a point of no return.",
        size=15,
        bold=True,
        color=WHITE,
    )

    # Right column — funnel
    _text(slide, Inches(7.0), Inches(1.75), Inches(5.5), Inches(0.35), "THE STAKES", size=11, bold=True, color=ACCENT)

    funnel_steps = [
        ("CLAIM FILE", "Notes · emails · calls · forms", BLUE),
        ("TODAY", "Manual review · slow · subjective", MUTED),
        ("ESCALATION OBVIOUS", "Window for cheap fix — closed", ACCENT),
        ("TOO LATE", "Legal fees ↑  ·  leverage ↓", RED),
    ]
    fy = 2.1
    for label, sub, col in funnel_steps:
        _box(slide, Inches(7.2), Inches(fy), Inches(5.4), Inches(0.72), RGBColor(30, 41, 59))
        _text(slide, Inches(7.4), Inches(fy + 0.08), Inches(5), Inches(0.3), label, size=13, bold=True, color=col)
        _text(slide, Inches(7.4), Inches(fy + 0.36), Inches(5), Inches(0.28), sub, size=11, color=MUTED)
        if label != "TOO LATE":
            _text(slide, Inches(9.5), Inches(fy + 0.78), Inches(0.5), Inches(0.25), "▼", size=14, color=MUTED, align=PP_ALIGN.CENTER)
        fy += 1.05

    # Stat callouts
    stats = [("$$$", "Legal spend can\nexceed claim value"), ("45 min", "Per case\nmanual research"), ("?", "Risk hidden in\nunstructured notes")]
    sx = 7.2
    for big, small in stats:
        _box(slide, Inches(sx), Inches(6.15), Inches(1.65), Inches(0.95), RGBColor(30, 41, 59))
        _text(slide, Inches(sx), Inches(6.2), Inches(1.65), Inches(0.45), big, size=22, bold=True, color=ACCENT, align=PP_ALIGN.CENTER)
        _text(slide, Inches(sx + 0.05), Inches(6.62), Inches(1.55), Inches(0.5), small, size=9, color=MUTED, align=PP_ALIGN.CENTER)
        sx += 1.85

    # Footer
    _text(
        slide,
        Inches(0.6),
        Inches(6.85),
        Inches(12),
        Inches(0.4),
        "RiskRadar — predict escalation early · explain why · act before legal costs run away",
        size=11,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
