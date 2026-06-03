"""Generate a single 16:9 future-scope slide (Google Slides import ready)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "RiskRadar_Future_Scope.pptx"

BG = RGBColor(15, 23, 42)       # slate-900
PANEL = RGBColor(30, 41, 59)    # slate-800
ACCENT = RGBColor(56, 189, 248) # sky-400
MUTED = RGBColor(148, 163, 184) # slate-400
WHITE = RGBColor(248, 250, 252)
GOOD = RGBColor(74, 222, 128)   # green-400


def _box(slide, left, top, width, height, fill=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def _text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    return tb


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _box(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, BG)
    _box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT)

    _text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.8), "Future Scope: From Decision Support to Decision Intelligence", size=32, bold=True)
    _text(slide, Inches(0.6), Inches(1.05), Inches(11.5), Inches(0.35), "Next 6-12 months roadmap for enterprise adoption", size=13, color=MUTED)

    items = [
        (
            "1) Closed-Loop Learning",
            "Capture adjuster outcomes and legal resolution status to continuously calibrate risk thresholds and improve recall/precision by segment.",
            "Outcome: measurable model lift with audit trail"
        ),
        (
            "2) Deeper Salesforce Integration",
            "Two-way sync: push RiskRadar scores, recommended actions, and follow-up tasks back to Salesforce so teams act inside existing workflows.",
            "Outcome: zero workflow disruption for claim teams"
        ),
        (
            "3) Governance and Trust Layer",
            "Add model cards, cohort fairness dashboards, drift alerts, and approval gates before high-impact recommendations are applied.",
            "Outcome: regulator and leadership confidence"
        ),
        (
            "4) Proactive Intervention Playbooks",
            "Auto-generate next-best-action playbooks by claim archetype (communication templates, escalation avoidance paths, owner SLAs).",
            "Outcome: faster time-to-intervention and lower leakage"
        ),
    ]

    y = 1.6
    for title, body, outcome in items:
        _box(slide, Inches(0.6), Inches(y), Inches(12.1), Inches(1.2), PANEL)
        _text(slide, Inches(0.85), Inches(y + 0.12), Inches(5.2), Inches(0.32), title, size=16, bold=True, color=ACCENT)
        _text(slide, Inches(0.85), Inches(y + 0.45), Inches(8.2), Inches(0.55), body, size=11, color=WHITE)
        _text(slide, Inches(9.2), Inches(y + 0.42), Inches(3.2), Inches(0.55), outcome, size=10, bold=True, color=GOOD, align=PP_ALIGN.RIGHT)
        y += 1.35

    _text(
        slide,
        Inches(0.6),
        Inches(6.95),
        Inches(12.1),
        Inches(0.3),
        "Vision: RiskRadar becomes the escalation intelligence layer across claims operations, legal strategy, and system-of-record workflows.",
        size=10,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
