"""Generate a single 16:9 tech-stack slide (Google Slides import ready)."""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parents[1]
OUT = ROOT / "docs" / "RiskRadar_Tech_Stack.pptx"

BG = RGBColor(15, 23, 42)
PANEL = RGBColor(30, 41, 59)
ACCENT = RGBColor(56, 189, 248)
MUTED = RGBColor(148, 163, 184)
WHITE = RGBColor(248, 250, 252)
TAG = RGBColor(96, 165, 250)


def _box(slide, left, top, width, height, fill=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    return shape


def _text(slide, left, top, width, height, text, size=18, bold=False, color=WHITE, align=PP_ALIGN.LEFT):
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = "Arial"
    return tb


def _panel(slide, left, top, width, height, title, bullets):
    _box(slide, left, top, width, height, PANEL)
    _text(slide, left + Inches(0.2), top + Inches(0.12), width - Inches(0.4), Inches(0.3), title, size=13, bold=True, color=ACCENT)
    body = "\n".join(f"• {b}" for b in bullets)
    _text(slide, left + Inches(0.2), top + Inches(0.42), width - Inches(0.35), height - Inches(0.5), body, size=10, color=WHITE)


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    _box(slide, Inches(0), Inches(0), prs.slide_width, prs.slide_height, BG)
    _box(slide, Inches(0), Inches(0), prs.slide_width, Inches(0.08), ACCENT)

    _text(slide, Inches(0.6), Inches(0.35), Inches(12), Inches(0.7), "Technology Stack", size=34, bold=True)
    _text(
        slide,
        Inches(0.6),
        Inches(1.0),
        Inches(12),
        Inches(0.35),
        "End-to-end escalation intelligence — from labelled claims to production adjuster workflow",
        size=12,
        color=MUTED,
    )

    col_w = Inches(3.95)
    col_h = Inches(2.55)
    gap_x = Inches(0.25)
    gap_y = Inches(0.22)
    x0 = Inches(0.55)
    y0 = Inches(1.45)

    panels = [
        (
            "Machine Learning & Data",
            [
                "Python 3.11 · Pandas · NumPy",
                "XGBoost (structured + text classifiers)",
                "scikit-learn (TF-IDF, Isotonic calibration, PSI)",
                "SHAP explainability · 29-feature matrix",
                "Time-to-escalation regressor (XGBoost)",
                "sentence-transformers (MiniLM-L6-v2, 550×384 index)",
            ],
        ),
        (
            "AI & Language Layer",
            [
                "Groq — Llama 3.1 / 3.3 (briefs, emails, actions)",
                "Google Gemini 2.5 Flash (second opinion)",
                "Temperature = 0 · citation-grounded prompts",
                "Trigger-phrase lexicon (rule-based NLP)",
                "RAG-style legal briefs with allowlist validation",
            ],
        ),
        (
            "Backend API",
            [
                "FastAPI + Uvicorn (REST)",
                "Pydantic · python-dotenv",
                "/predict · /claims · /metrics · /feedback",
                "Risk calibrator (64/28/8 blend)",
                "Similar-claims NN · drift · fairness endpoints",
            ],
        ),
        (
            "Frontend",
            [
                "Next.js 16 · React 19 · TypeScript",
                "Tailwind CSS 4",
                "Triage queue · single-claim deep dive",
                "Live ROI calculator · adjuster verdict UI",
                "ReactMarkdown for AI brief rendering",
            ],
        ),
        (
            "Integrations & Legal Data",
            [
                "Salesforce — Apex Case sync + REST webhook",
                "Indian Kanoon API (India precedents)",
                "CourtListener API (US case law)",
                "External claims store (SF + demo cases)",
            ],
        ),
        (
            "Deploy & Persistence",
            [
                "Docker (FastAPI + Next.js + nginx)",
                "Render (cloud web services)",
                "SQLite — adjuster feedback log",
                "joblib / JSON — trained model artifacts",
                "Excel/CSV — 550-claim labelled dataset",
            ],
        ),
    ]

    for i, (title, bullets) in enumerate(panels):
        row, col = divmod(i, 3)
        x = x0 + col * (col_w + gap_x)
        y = y0 + row * (col_h + gap_y)
        _panel(slide, x, y, col_w, col_h, title, bullets)

    _text(
        slide,
        Inches(0.6),
        Inches(6.85),
        Inches(12.1),
        Inches(0.45),
        "Core principle: ML scores the risk · LLMs explain and draft · humans decide · every legal citation is retrieved and validated",
        size=10,
        color=TAG,
        align=PP_ALIGN.CENTER,
    )

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUT)
    print(f"Wrote {OUT}")


if __name__ == "__main__":
    build()
